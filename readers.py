"""Text extraction from uploaded documents."""

import os

import html2text
import pandas as pd
from docx import Document
from pptx import Presentation
from pypdf import PdfReader

from log import logger

# A PDF that yields less than this many characters per page is almost certainly
# scanned images rather than text.
SCANNED_PDF_CHARS_PER_PAGE = 32


class EmptyDocumentError(Exception):
    """Raised when a document parsed fine but produced no usable text."""


class ScannedPdfError(EmptyDocumentError):
    """Raised when a PDF looks like page images with no extractable text layer."""


def read_pdf_file(file_path):
    text_parts = []
    with open(file_path, "rb") as file:
        reader = PdfReader(file)
        page_count = len(reader.pages)
        for page in reader.pages:
            text_parts.append(page.extract_text() or "")

    text = "\n".join(text_parts)
    if page_count and len(text.strip()) < SCANNED_PDF_CHARS_PER_PAGE * page_count:
        raise ScannedPdfError(
            f"This PDF has {page_count} page(s) but almost no extractable text, so it is "
            "probably scanned images. Select a model that uses the Responses API to have "
            "the PDF read natively, or run OCR on it first."
        )
    return text


def read_word_file(file_path):
    document = Document(file_path)
    return "\n".join(paragraph.text for paragraph in document.paragraphs)


def read_excel_file(file_path):
    """Read every sheet, not just the first one."""
    sheets = pd.read_excel(file_path, sheet_name=None)
    parts = []
    for sheet_name, frame in sheets.items():
        parts.append(f"# Sheet: {sheet_name}\n{frame.to_csv(index=False)}")
    return "\n\n".join(parts)


def read_powerpoint_file(file_path):
    presentation = Presentation(file_path)
    parts = []
    for index, slide in enumerate(presentation.slides, start=1):
        slide_text = [shape.text for shape in slide.shapes if hasattr(shape, "text")]
        slide_text = [line for line in slide_text if line.strip()]
        if slide_text:
            parts.append(f"# Slide {index}\n" + "\n".join(slide_text))
    return "\n\n".join(parts)


def read_html_file(file_path):
    return html2text.html2text(read_text_file(file_path))


def read_text_file(file_path):
    encodings = ["utf-8", "utf-8-sig", "cp1252", "latin-1"]
    for encoding in encodings:
        try:
            with open(file_path, "r", encoding=encoding) as file:
                return file.read()
        except UnicodeDecodeError:
            continue
    raise UnicodeDecodeError(
        "text", b"", 0, 1, f"could not decode as any of: {', '.join(encodings)}"
    )


READERS_BY_EXTENSION = {
    ".pdf": read_pdf_file,
    ".docx": read_word_file,
    ".xlsx": read_excel_file,
    ".xls": read_excel_file,
    ".pptx": read_powerpoint_file,
    ".htm": read_html_file,
    ".html": read_html_file,
}


def read_document(file_path):
    """Extract text from a document, dispatching on file extension.

    Raises EmptyDocumentError when the file parsed but yielded nothing usable, and
    lets the underlying parser's own exception through otherwise.
    """
    _, extension = os.path.splitext(file_path)
    reader = READERS_BY_EXTENSION.get(extension.lower(), read_text_file)
    logger.debug("Reading %s with %s", os.path.basename(file_path), reader.__name__)

    text = reader(file_path)
    if not text or not text.strip():
        raise EmptyDocumentError("No text could be extracted from this file.")
    return text
