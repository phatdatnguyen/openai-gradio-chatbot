import os
import io
import json
import glob
import re
import requests
import tiktoken
import html2text
from PIL import Image
import base64
import gradio as gr
from openai import OpenAI
import PyPDF2
from docx import Document
import pandas as pd
from pptx import Presentation

try:
    from api_key import API_KEY as FILE_API_KEY
except ImportError:
    FILE_API_KEY = None


API_KEY = os.getenv("OPENAI_API_KEY") or FILE_API_KEY
client = OpenAI(api_key=API_KEY) if API_KEY else OpenAI()
IMAGE_DATA_URL_PREFIX = "data:image/png;base64,"


def normalize_text(value):
    if value is None:
        return ""
    return str(value).strip()


def serialize_message_value(value):
    if value is None:
        return ""
    if isinstance(value, str):
        return value
    return json.dumps(value, ensure_ascii=False)

def encode_image(image):
    buffered = io.BytesIO()
    image.save(buffered, format="PNG")
    return base64.b64encode(buffered.getvalue()).decode('utf-8')

def encode_image_from_url(url):
    response = requests.get(url, timeout=15)
    response.raise_for_status()
    image = Image.open(io.BytesIO(response.content))
    return encode_image(image)

def prepare_chat_messages(history):
    messages = []
    for message in history or []:
        content = message.get("content", "")
        if message["role"] == "user" and message.get("image_url"):
            user_content = []
            if content:
                user_content.append({"type": "text", "text": content})
            user_content.append({
                "type": "image_url",
                "image_url": {"url": message["image_url"]}
            })
            messages.append({"role": "user", "content": user_content})
        else:
            messages.append({"role": message["role"], "content": content})
    return messages


def prepare_responses_input(history):
    messages = []
    for message in history or []:
        content = message.get("content", "")
        if message["role"] == "user" and message.get("image_url"):
            user_content = []
            if content:
                user_content.append({"type": "input_text", "text": content})
            user_content.append({"type": "input_image", "image_url": message["image_url"]})
            messages.append({"role": "user", "content": user_content})
        else:
            messages.append({"role": message["role"], "content": content})
    return messages


def process_image(llm_model, temperature, top_p, text, image_path, history, system_prompt=""):
    prompt_text = normalize_text(text)
    image = Image.open(image_path)
    base64_image = encode_image(image)
    image_data_url = f"{IMAGE_DATA_URL_PREFIX}{base64_image}"

    history = history or []
    history.append({"role": "user", "content": prompt_text, "image_url": image_data_url})

    print(f'User: {prompt_text}\nImage: {image}')

    chat_messages = trim_history(prepare_chat_messages(history), llm_model)
    if system_prompt:
        chat_messages = [{"role": "system", "content": system_prompt}] + chat_messages
    response = client.chat.completions.create(
        model=llm_model,
        messages=chat_messages,
        temperature=temperature,
        top_p=top_p)

    ai_message = response.choices[0].message.content.strip()
    history.append({"role": "assistant", "content": ai_message})

    print(f'AI: {ai_message}')

    return history

def read_pdf_file(file_path):
    text = ""
    with open(file_path, 'rb') as file:
        reader = PyPDF2.PdfReader(file)
        for page in reader.pages:
            text += page.extract_text() or ""
            text += '\n'
    return text

def read_word_file(file_path):
    document = Document(file_path)
    text = []
    for paragraph in document.paragraphs:
        text.append(paragraph.text)
    return '\n'.join(text)

def read_excel_file(file_path):
    excel_data = pd.read_excel(file_path)
    return excel_data.to_csv(index=False)

def read_powerpoint_file(file_path):
    presentation = Presentation(file_path)
    text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return '\n'.join(text)

def read_html_file(file_path):
    with open(file_path, 'tr', encoding="utf-8") as file:
        text = file.read()
    return html2text.html2text(text)

def read_text_file(file_path):
    encodings = ["utf-8", "utf-8-sig", "cp1252", "latin-1"]
    for encoding in encodings:
        try:
            with open(file_path, 'tr', encoding=encoding) as file:
                return file.read()
        except UnicodeDecodeError:
            continue
    raise UnicodeDecodeError("text", b"", 0, 1, "Unsupported text encoding")

def replace_history_content(history):
    replaced_history = []
    for message in history:
        content = message["content"]

        pattern = r'(?s)<<<DOCUMENT_CONTENT>>>\nFile: (.*?)\n.*?<<<END_DOCUMENT>>>'
        def repl(match):
            file_name = match.group(1)
            return f'[File: {file_name}]'
        replaced_content = re.sub(pattern, repl, content)

        pattern = r'(?s)<<<LINK_CONTENT>>>\nURL: (.*?)\n.*?<<<END_LINK>>>'
        def repl(match):
            url = match.group(1)
            return f'[URL: {url}]'
        replaced_content = re.sub(pattern, repl, replaced_content)

        if "image_url" in message:
            image_url = message["image_url"]
            replaced_content += f'\n![Image]({image_url})'

        replaced_message = {
                "role": message["role"],
                "content": replaced_content
            }
        replaced_history.append(replaced_message)

    return replaced_history

MODEL_TOKEN_LIMITS = {
    "gpt-3.5-turbo": 4096,
    "gpt-4": 8192,
    "gpt-4-turbo": 128000,
    "gpt-4.1": 1047576,
    "gpt-4.1-mini": 1047576,
    "gpt-4.1-nano": 1047576,
    "gpt-4o": 128000,
    "gpt-4o-mini": 128000,
    "gpt-5": 400000,
    "gpt-5-chat-latest": 128000,
    "gpt-5-pro": 400000,
    "gpt-5-mini": 400000,
    "gpt-5-nano": 400000,
    "gpt-5.1": 400000,
    "gpt-5.1-codex": 400000,
    "gpt-5.1-codex-max": 400000,
    "gpt-5.1-codex-mini": 400000,
    "gpt-5.2": 400000,
    "gpt-5.2-codex": 400000,
    "gpt-5.2-pro": 400000,
    "gpt-5.3-codex": 400000,
    "gpt-5.3-chat": 400000,
    "gpt-5.4": 1050000,
    "gpt-5.4-nano": 400000,
    "gpt-5.4-mini": 400000,
    "gpt-5.4-pro": 1050000,
    "o1": 128000,
    "o1-mini": 128000,
    "o1-pro": 128000,
    "o3": 200000,
    "o3-mini": 200000,
    "o3-pro": 200000,
    "o3-deep-research": 200000,
    "o4-mini": 200000,
    "o4-mini-deep-research": 200000
}

MODEL_TOKEN_LIMITS_WITH_WEB_SEARCH = {
    "gpt-4.1": 128000,
    "gpt-4.1-mini": 128000,
    "gpt-4o": 128000,
    "gpt-4o-mini": 128000,
    "gpt-5": 400000,
    "gpt-5-chat-latest": 128000,
    "gpt-5-pro": 400000,
    "gpt-5-mini": 400000,
    "gpt-5.1": 400000,
    "gpt-5.2": 400000,
    "gpt-5.2-pro": 400000,
    "gpt-5.4": 1050000,
    "gpt-5.4-nano": 400000,
    "gpt-5.4-mini": 400000,
    "gpt-5.4-pro": 1050000,
    "o1": 128000,
    "o3": 128000,
    "o3-pro": 128000,
    "o3-deep-research": 128000,
    "o4-mini-deep-research": 128000
}

MODEL_MAX_OUTPUT_TOKENS = {
    "gpt-5.1-codex-max": 128000,
    "gpt-5.2-codex": 128000,
    "gpt-5.3-codex": 128000,
    "gpt-5.3-chat": 128000,
    "gpt-5.4": 128000,
    "gpt-5.4-nano": 128000,
    "gpt-5.4-mini": 128000,
    "gpt-5.4-pro": 128000
}

RESPONSES_API_MODELS = {
    "o1-pro",
    "o3-pro",
    "o3-deep-research",
    "o4-mini-deep-research",
    "gpt-5.1-codex",
    "gpt-5.1-codex-max",
    "gpt-5.1-codex-mini",
    "gpt-5.2-codex",
    "gpt-5.3-codex"
}

WEB_SEARCH_MODELS = {
    "gpt-4.1",
    "gpt-4.1-mini",
    "gpt-4o",
    "gpt-4o-mini",
    "gpt-5",
    "gpt-5-chat-latest",
    "gpt-5-pro",
    "gpt-5-mini",
    "gpt-5.1",
    "gpt-5.2",
    "gpt-5.2-pro",
    "gpt-5.4",
    "gpt-5.4-nano",
    "gpt-5.4-mini",
    "gpt-5.4-pro",
    "o3",
    "o3-pro"
}

VISION_DISABLED_MODELS = {"gpt-3.5-turbo", "o1-mini", "o3-mini"}

IMAGE_MODEL_CONFIGS = {
    "gpt-image-2": {
        "size_choices": ["auto", "1024x1024", "1536x1024", "1024x1536", "2048x2048", "2048x1152", "3840x2160", "2160x3840"],
        "default_size": "auto",
        "quality_choices": ["auto", "low", "medium", "high"],
        "default_quality": "auto",
        "background_choices": ["auto", "opaque"],
        "default_background": "auto",
        "input_fidelity_choices": ["high"],
        "default_input_fidelity": "high",
        "input_fidelity_interactive": False,
    },
    "gpt-image-1.5": {
        "size_choices": ["auto", "1024x1024", "1536x1024", "1024x1536"],
        "default_size": "auto",
        "quality_choices": ["auto", "low", "medium", "high"],
        "default_quality": "auto",
        "background_choices": ["auto", "opaque", "transparent"],
        "default_background": "auto",
        "input_fidelity_choices": ["low", "high"],
        "default_input_fidelity": "high",
        "input_fidelity_interactive": True,
    },
}

IMAGE_MODEL_CHOICES = list(IMAGE_MODEL_CONFIGS.keys())
IMAGE_OUTPUT_FORMAT_CHOICES = ["png", "jpeg", "webp"]
IMAGE_MODERATION_CHOICES = ["auto", "low"]

MODEL_CHOICES = [
    "gpt-3.5-turbo",
    "gpt-4",
    "gpt-4-turbo",
    "gpt-4.1",
    "gpt-4.1-mini",
    "gpt-4.1-nano",
    "gpt-4o",
    "gpt-4o-mini",
    "gpt-5",
    "gpt-5-chat-latest",
    "gpt-5-pro",
    "gpt-5-mini",
    "gpt-5-nano",
    "gpt-5.1",
    "gpt-5.1-codex",
    "gpt-5.1-codex-max",
    "gpt-5.1-codex-mini",
    "gpt-5.2",
    "gpt-5.2-codex",
    "gpt-5.2-pro",
    "gpt-5.3-codex",
    "gpt-5.3-chat",
    "gpt-5.4",
    "gpt-5.4-nano",
    "gpt-5.4-mini",
    "gpt-5.4-pro",
    "o1",
    "o1-mini",
    "o1-pro",
    "o3",
    "o3-mini",
    "o3-pro",
    "o3-deep-research",
    "o4-mini",
    "o4-mini-deep-research"
]


def get_image_model_config(image_model):
    return IMAGE_MODEL_CONFIGS.get(image_model, IMAGE_MODEL_CONFIGS["gpt-image-2"])


def build_image_option_updates(image_model):
    config = get_image_model_config(image_model)
    return (
        gr.Dropdown(label="Image size", value=config["default_size"], choices=config["size_choices"], interactive=True),
        gr.Dropdown(label="Image quality", value=config["default_quality"], choices=config["quality_choices"], interactive=True),
        gr.Dropdown(label="Image background", value=config["default_background"], choices=config["background_choices"], interactive=True),
        gr.Dropdown(
            label="Input fidelity",
            value=config["default_input_fidelity"],
            choices=config["input_fidelity_choices"],
            interactive=config["input_fidelity_interactive"],
        ),
    )

def on_image_output_format_change(image_output_format):
    compression_enabled = image_output_format in ["jpeg", "webp"]
    return gr.Slider(label="Compression", minimum=0, maximum=100, step=1, value=100, interactive=compression_enabled)

def get_max_context_tokens(model_name, web_search=False):
    if web_search:
        return MODEL_TOKEN_LIMITS_WITH_WEB_SEARCH.get(model_name, 128000)
    else:
        return MODEL_TOKEN_LIMITS.get(model_name, 4096)

def count_tokens(messages, model):
    try:
        encoding = tiktoken.encoding_for_model(model)
    except:
        encoding = tiktoken.get_encoding("cl100k_base")
    num_tokens = 0
    for msg in messages:
        num_tokens += 4
        for key, value in msg.items():
            num_tokens += len(encoding.encode(serialize_message_value(value)))
    num_tokens += 2
    return num_tokens

def trim_history(history, model, web_search=False, reserved_tokens=2000):
    if not history:
        return []

    trimmed_history = []
    total_tokens = 0
    max_context = get_max_context_tokens(model, web_search)
    max_input_tokens = max_context - reserved_tokens

    for message in reversed(history):
        message_tokens = count_tokens([message], model)
        if total_tokens + message_tokens <= max_input_tokens:
            trimmed_history.insert(0, message)
            total_tokens += message_tokens
        else:
            break

    return trimmed_history

def process_image_generation(text, image_path, history, image_model, image_size, image_quality, image_background, image_output_format, image_output_compression, image_moderation, image_input_fidelity):
    history = history or []
    prompt_text = normalize_text(text)
    if not prompt_text:
        gr.Warning("Enter a prompt before generating or editing an image.")
        return history

    user_message = {"role": "user", "content": prompt_text}
    if image_path:
        with Image.open(image_path) as input_image:
            user_message["image_url"] = f"{IMAGE_DATA_URL_PREFIX}{encode_image(input_image)}"
    history.append(user_message)

    print(f'User: {prompt_text}')

    image_kwargs = {
        "model": image_model,
        "prompt": prompt_text,
        "size": image_size,
        "quality": image_quality,
        "background": image_background,
        "output_format": image_output_format,
    }
    if image_output_format in ["jpeg", "webp"]:
        image_kwargs["output_compression"] = int(image_output_compression)

    if image_path:
        if image_model == "gpt-image-2":
            image_kwargs.pop("response_format", None)
            image_kwargs.pop("input_fidelity", None)
        else:
            image_kwargs["input_fidelity"] = image_input_fidelity

        with open(image_path, "rb") as input_image_file:
            response = client.images.edit(image=input_image_file, **image_kwargs)
    else:
        image_kwargs["moderation"] = image_moderation
        response = client.images.generate(**image_kwargs)

    image_response = response.data[0] if response.data else None
    base64_image = getattr(image_response, "b64_json", None) if image_response else None
    revised_prompt = getattr(image_response, "revised_prompt", None) if image_response else None

    if base64_image:
        action_label = "Edited image" if image_path else "Generated image"
        ai_message = f"{action_label} with `{image_model}`."
        if revised_prompt and revised_prompt != prompt_text:
            ai_message += f"\n\nRevised prompt:\n{revised_prompt}"
        ai_message += f"\n\n![Image]({IMAGE_DATA_URL_PREFIX}{base64_image})"
        history.append({"role": "assistant", "content": ai_message})

    return history

def process_document(llm_model, web_search, temperature, top_p, text, document, history, system_prompt=""):
    document_text = ""
    prompt_text = normalize_text(text)
    try:
        _, file_extension = os.path.splitext(document)
        file_name = os.path.basename(document)
        if file_extension.lower() == ".pdf":
            document_text = read_pdf_file(document)
        elif file_extension.lower() == ".docx":
            document_text = read_word_file(document)
        elif file_extension.lower() == ".xlsx":
            document_text = read_excel_file(document)
        elif file_extension.lower() == ".pptx":
            document_text = read_powerpoint_file(document)
        elif file_extension.lower() in (".htm", ".html"):
            document_text = read_html_file(document)
        else:
            document_text = read_text_file(document)

        document_text = f"<<<DOCUMENT_CONTENT>>>\nFile: {file_name}\n{document_text}\n<<<END_DOCUMENT>>>"
    except Exception as exc:
        gr.Warning("Cannot read this file!\n" + str(exc))
        return history or []

    if document_text:
        history = history or []
        parts = [part for part in [prompt_text, document_text] if part]
        history.append({"role": "user", "content": '\n'.join(parts)})

        print(f'User: {prompt_text}\nUser uploaded file: {os.path.basename(document)}')
    else:
        gr.Warning("The uploaded document appears to be empty.")
        return history or []

    responses_kwargs = {}
    if system_prompt:
        responses_kwargs["instructions"] = system_prompt

    if web_search != "None":
        response = client.responses.create(
            model=llm_model,
            input=trim_history(prepare_responses_input(history), llm_model, web_search=True),
            tools=[{"type": "web_search_preview", "search_context_size": web_search}],
            temperature=temperature,
            top_p=top_p,
            **responses_kwargs,
        )
        ai_message = response.output_text
    elif llm_model in RESPONSES_API_MODELS:
        response = client.responses.create(
            model=llm_model,
            input=trim_history(prepare_responses_input(history), llm_model),
            temperature=temperature,
            top_p=top_p,
            **responses_kwargs,
        )
        ai_message = response.output_text
    else:
        messages = trim_history(prepare_chat_messages(history), llm_model)
        if system_prompt:
            messages = [{"role": "system", "content": system_prompt}] + messages
        response = client.chat.completions.create(
            model=llm_model,
            messages=messages,
            temperature=temperature,
            top_p=top_p,
        )
        ai_message = response.choices[0].message.content.strip()

    history.append({"role": "assistant", "content": ai_message})

    print(f'AI: {ai_message}')

    return history

def process_text(llm_model, web_search, temperature, top_p, text, url, history, system_prompt=""):
    history = history or []
    prompt_text = normalize_text(text)
    link_text = ""

    if url:
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/114.0.0.0 Safari/537.36"
        }
        try:
            response = requests.get(url, headers=headers, timeout=15)
            response.raise_for_status()
            response.encoding = response.encoding or "utf-8"
            html_text = html2text.html2text(response.text)
            link_text = f'<<<LINK_CONTENT>>>\nURL: {url}\n{html_text}\n<<<END_LINK>>>'
        except requests.RequestException as exc:
            gr.Warning(f"Failed to retrieve the webpage: {exc}")

    if not prompt_text and not link_text:
        gr.Warning("Enter a message, add a link, or upload a file/image before sending.")
        return history

    combined_text = '\n'.join(part for part in [prompt_text, link_text] if part)
    history.append({"role": "user", "content": combined_text})

    print(f'User: {combined_text}')

    responses_kwargs = {}
    if system_prompt:
        responses_kwargs["instructions"] = system_prompt

    if web_search != "None":
        response = client.responses.create(
            model=llm_model,
            input=trim_history(prepare_responses_input(history), llm_model, web_search=True),
            tools=[{"type": "web_search_preview", "search_context_size": web_search}],
            temperature=temperature,
            top_p=top_p,
            **responses_kwargs,
        )
        ai_message = response.output_text
    elif llm_model in RESPONSES_API_MODELS:
        response = client.responses.create(
            model=llm_model,
            input=trim_history(prepare_responses_input(history), llm_model),
            temperature=temperature,
            top_p=top_p,
            **responses_kwargs,
        )
        ai_message = response.output_text
    else:
        messages = trim_history(prepare_chat_messages(history), llm_model)
        if system_prompt:
            messages = [{"role": "system", "content": system_prompt}] + messages
        response = client.chat.completions.create(
            model=llm_model,
            messages=messages,
            temperature=temperature,
            top_p=top_p,
        )
        ai_message = response.choices[0].message.content.strip()

    history.append({"role": "assistant", "content": ai_message})

    print(f'AI: {ai_message}')

    return history

def on_llm_model_change(llm_model):
    if llm_model in ["o3-deep-research", "o4-mini-deep-research"]:
       web_search = gr.Dropdown(label="Web search", value="medium", choices=["low", "medium", "high"], interactive=True)
    elif llm_model in WEB_SEARCH_MODELS:
       web_search = gr.Dropdown(label="Web search", value="None", choices=["None", "low", "medium", "high"], interactive=True)
    else:
       web_search = gr.Dropdown(label="Web search", value="None", choices=["None", "low", "medium", "high"], interactive=False)

    image_input = gr.Image(label="Upload an image", sources=["upload", "clipboard"], type="filepath", interactive=True)
    generate_image = gr.Checkbox(label="Generate or edit image", value=False, interactive=True)

    return web_search, image_input, generate_image

def make_reset_inputs():
    return (
        gr.Textbox(label="Message", placeholder="Type a message or question...", autofocus=True, value=None),
        gr.Image(label="Upload an image", sources=["upload", "clipboard"], type="filepath", value=None),
        gr.File(label="Upload a document", type="filepath", value=None),
        gr.Textbox(label="Link", value=None),
        gr.Checkbox(label="Generate or edit image", value=False),
    )

def on_user_input(llm_model, web_search, temperature, top_p, text, image, document, url, history, generate_image, system_prompt, image_model, image_size, image_quality, image_background, image_output_format, image_output_compression, image_moderation, image_input_fidelity):
    history = history or []
    try:
        if generate_image:
            history = process_image_generation(
                text, image, history, image_model, image_size, image_quality,
                image_background, image_output_format, image_output_compression,
                image_moderation, image_input_fidelity,
            )
        elif image:
            if llm_model in VISION_DISABLED_MODELS:
                gr.Warning("The selected chat model does not support image analysis. Enable image generation to edit the image, or choose a vision-capable chat model.")
            else:
                history = process_image(llm_model, temperature, top_p, text, image, history, system_prompt)
        elif document:
            history = process_document(llm_model, web_search, temperature, top_p, text, document, history, system_prompt)
        elif text or url:
            history = process_text(llm_model, web_search, temperature, top_p, text, url, history, system_prompt)
    except Exception as exc:
        gr.Warning(str(exc))

    replaced_history = replace_history_content(history)
    return (history, replaced_history) + make_reset_inputs()

def on_new_chat_click():
    return [], []

def on_toggle_history_column(state):
    state = not state
    return gr.update(visible=state), state

def get_history_file_list():
    history_file_list = []
    for file in glob.glob('./history/*.json', recursive=True):
        history_file_list.append(os.path.splitext(os.path.basename(file))[0])

    return sorted(history_file_list, key=str.lower)

def select_history_file(evt: gr.SelectData):
    row_index = evt.index[0]
    history_files = get_history_file_list()
    if row_index >= len(history_files):
        return ""
    history_file = history_files[row_index]
    return history_file

def sanitize_history_file_name(history_file_name):
    if hasattr(history_file_name, "value"):
        history_file_name = history_file_name.value
    sanitized_name = re.sub(r'[<>:"/\\|?*\x00-\x1f]', "_", normalize_text(history_file_name))
    return sanitized_name[:100] or "Chat history"

def save_history(history, history_file_name):
    try:
        if not history:
            raise Exception("No history")

        os.makedirs("history", exist_ok=True)
        safe_history_file_name = sanitize_history_file_name(history_file_name)
        with open(f"./history/{safe_history_file_name}.json", "w", encoding='utf8') as file:
            json.dump(history, file, indent=4, ensure_ascii=False)
        return "Chat history saved successfully!", pd.DataFrame(get_history_file_list(), columns=["File name"])
    except Exception as exc:
        return f"Error saving history: {str(exc)}", pd.DataFrame(get_history_file_list(), columns=["File name"])

def load_history(history_file_name):
    try:
        safe_history_file_name = sanitize_history_file_name(history_file_name)
        history_file_path = f"./history/{safe_history_file_name}.json"
        with open(history_file_path, "r", encoding='utf8') as file:
            history = json.load(file)
        return history, replace_history_content(history), "Chat history loaded successfully!"
    except Exception as exc:
        return [], [], f"Error loading history: {str(exc)}"

with gr.Blocks() as demo:
    with gr.Row(equal_height=True):
        with gr.Column(scale=4):
            chatbot = gr.Chatbot(buttons=["copy"], min_height=800)
            state = gr.State([])
            history_column_state = gr.State(True)
            with gr.Row(equal_height=True):
                new_chat_button = gr.Button(value="New chat")
                toggle_history_column_button = gr.Button(value="Toggle chat history")
        with gr.Column(scale=1) as history_column:
            history_files_list = gr.DataFrame(label="Chat history files", value=pd.DataFrame(get_history_file_list(), columns=["File name"]), max_height=260)
            history_file_name = gr.Textbox(label="File name", value="Chat history")
            load_button = gr.Button(value="Load chat history")
            load_status = gr.Markdown(value="")
            save_button = gr.Button(value="Save chat history")
            save_status = gr.Markdown(value="")
    with gr.Row():
        with gr.Column(scale=1):
            with gr.Accordion(label="Prompt"):
                with gr.Row(equal_height=True):
                    text_input = gr.Textbox(label="Message", placeholder="Type a message or question...", autofocus=True, scale=4)
                    send_button = gr.Button(value="Send", variant="primary", scale=1, min_width=80)
                system_prompt = gr.Textbox(label="System prompt", placeholder="Optional instructions for the AI...", lines=2)
                llm_model = gr.Dropdown(label="Model", value="gpt-5.4", choices=MODEL_CHOICES)
                web_search = gr.Dropdown(label="Web search", value="None", choices=["None", "low", "medium", "high"])
                temperature = gr.Slider(label="Temperature", minimum=0, maximum=2, step=0.01, value=1)
                top_p = gr.Slider(label="Top-p", minimum=0, maximum=1, step=0.01, value=1)
        with gr.Column(scale=1):
            with gr.Accordion(label="Input"):
                image_input = gr.Image(label="Upload an image", sources=["upload", "clipboard"], type="filepath")
                document_input = gr.File(label="Upload a document", type="filepath")
                url_input = gr.Textbox(label="Link")
        with gr.Column(scale=1):
            with gr.Accordion(label="Image generation and editing"):
                default_image_model = "gpt-image-2"
                default_image_config = get_image_model_config(default_image_model)
                generate_image = gr.Checkbox(label="Generate or edit image", value=False)
                image_model = gr.Dropdown(label="Image model", value=default_image_model, choices=IMAGE_MODEL_CHOICES)
                image_size = gr.Dropdown(label="Image size", value=default_image_config["default_size"], choices=default_image_config["size_choices"])
                image_quality = gr.Dropdown(label="Image quality", value=default_image_config["default_quality"], choices=default_image_config["quality_choices"])
                image_background = gr.Dropdown(label="Image background", value=default_image_config["default_background"], choices=default_image_config["background_choices"])
                image_output_format = gr.Dropdown(label="Output format", value="png", choices=IMAGE_OUTPUT_FORMAT_CHOICES)
                image_output_compression = gr.Slider(label="Compression", minimum=0, maximum=100, step=1, value=100, interactive=False)
                image_moderation = gr.Dropdown(label="Moderation", value="auto", choices=IMAGE_MODERATION_CHOICES)
                image_input_fidelity = gr.Dropdown(
                    label="Input fidelity",
                    value=default_image_config["default_input_fidelity"],
                    choices=default_image_config["input_fidelity_choices"],
                    interactive=default_image_config["input_fidelity_interactive"],
                )

    user_input_inputs = [llm_model, web_search, temperature, top_p, text_input, image_input, document_input, url_input, state, generate_image, system_prompt, image_model, image_size, image_quality, image_background, image_output_format, image_output_compression, image_moderation, image_input_fidelity]
    user_input_outputs = [state, chatbot, text_input, image_input, document_input, url_input, generate_image]

    new_chat_button.click(on_new_chat_click, [], [chatbot, state])
    toggle_history_column_button.click(on_toggle_history_column, history_column_state, [history_column, history_column_state])

    llm_model.change(on_llm_model_change, llm_model, [web_search, image_input, generate_image])
    image_model.change(build_image_option_updates, image_model, [image_size, image_quality, image_background, image_input_fidelity])
    image_output_format.change(on_image_output_format_change, image_output_format, image_output_compression)
    text_input.submit(on_user_input, user_input_inputs, user_input_outputs)
    send_button.click(on_user_input, user_input_inputs, user_input_outputs)

    history_files_list.select(select_history_file, [], history_file_name)
    load_button.click(load_history, history_file_name, [state, chatbot, load_status])
    save_button.click(save_history, [state, history_file_name], [save_status, history_files_list])

if __name__ == "__main__":
    demo.launch(max_file_size=100*gr.FileSize.MB)
